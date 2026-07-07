"""Inspect the duotone template to understand layouts, fonts, colors."""
import os
from pptx import Presentation
from pptx.util import Emu

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(ROOT, "24_Template_PowerPoint_Duotone.pptx")


def main():
    prs = Presentation(TEMPLATE)
    print(f"Slide width:  {prs.slide_width}  ({prs.slide_width/914400:.2f} in)")
    print(f"Slide height: {prs.slide_height}  ({prs.slide_height/914400:.2f} in)")
    print(f"\nLAYOUTS ({len(prs.slide_layouts)}):")
    for i, lay in enumerate(prs.slide_layouts):
        print(f"  [{i}] {lay.name}  ({len(lay.placeholders)} placeholders)")
        for ph in lay.placeholders:
            print(f"        ph idx={ph.placeholder_format.idx:>3}  "
                  f"type={ph.placeholder_format.type}  name={ph.name}")
    print(f"\nMASTERS ({len(prs.slide_masters)}):")
    for sm in prs.slide_masters:
        print(f"  master, {len(sm.slide_layouts)} layouts")
    print(f"\nEXISTING SLIDES IN TEMPLATE ({len(prs.slides)}):")
    for i, s in enumerate(prs.slides):
        layout = s.slide_layout.name if s.slide_layout else "?"
        print(f"  Slide {i}: layout={layout}, {len(s.shapes)} shapes")
        for sh in s.shapes:
            txt = ""
            if sh.has_text_frame:
                txt = sh.text_frame.text[:60].replace("\n", " | ")
            print(f"        - {sh.shape_type}  name={sh.name}  text={txt!r}")


if __name__ == "__main__":
    main()
