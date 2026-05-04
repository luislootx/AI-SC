"""ICERM poster builder - 48"x36" landscape, layout v3.

Layout plan
-----------
  TOP BANNER (full width, h=4.6")
    - Top-left 4.5"x4.6" reserved for TAMU logo (user adds after build)
    - Title + authors + workshop info centered in the remaining width
    - Top-right ~4.0"x4.0" reserved for repo QR code

  BODY (3 logical zones):
    A. LEFT column (1/3 width)   - prose: Key Takeaway, Motivation, Methods,
                                   Discussion, References
    B. RIGHT 2-col SPAN (2/3 width):
         - Framework Overview (full span width, ~12")
         - Results - All Four Experiments (full span width, ~14")
         - Bottom row: Discovered Architectures + Honest Comparison side
           by side (~5")

Output: draft/poster_ICERM.pptx
"""
from __future__ import annotations
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

POSTER_DIR = r"C:/Users/luisl/OneDrive/Documentos/TAMU/2026.4 ICERM/Poster"
FIG_DIR = os.path.join(POSTER_DIR, "figures")
OUT = os.path.join(POSTER_DIR, "draft", "poster_ICERM.pptx")

# Colors
MAROON = RGBColor(0x50, 0x00, 0x00)
MAROON_DK = RGBColor(0x3F, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF4)
SOFT_BG = RGBColor(0xFD, 0xFD, 0xF6)

# ---------------------------------------------------------------- helpers

def _set_text(text_frame, text, *, font="Calibri", size=18, bold=False,
              color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              clear=True):
    if clear:
        text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.alignment = align
        p.text = line
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def _add_rect(slide, x, y, w, h, *, fill=MAROON, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    sh.shadow.inherit = False
    return sh


def _add_text_box(slide, x, y, w, h, text, **kwargs):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.margin_left = Inches(0.10)
    tb.text_frame.margin_right = Inches(0.10)
    tb.text_frame.margin_top = Inches(0.05)
    tb.text_frame.margin_bottom = Inches(0.05)
    _set_text(tb.text_frame, text, **kwargs)
    return tb


def estimate_body_inches(text, size_pt, width_emu):
    """Crude line-height calculation. Slightly generous to avoid cutoff."""
    width_in = width_emu / Inches(1)
    char_w_in = size_pt / 130.0   # tighter than 144 for Calibri
    chars_per_line = max(1, int(width_in / char_w_in))
    line_h_in = size_pt * 1.40 / 72.0
    n_lines = 0
    for line in text.split("\n"):
        if line == "":
            n_lines += 0.5
        else:
            n_lines += max(1, (len(line) + chars_per_line - 1) // chars_per_line)
    return max(0.7, n_lines * line_h_in + 0.40)


def add_section(slide, x, y, w, header, body_text, *, body_size=24,
                header_size=32, body_height=None):
    HDR_H = Inches(1.10)
    _add_rect(slide, x, y, w, HDR_H, fill=MAROON)
    _add_text_box(slide, x, y, w, HDR_H, header,
                  font="Calibri", size=header_size, bold=True,
                  color=WHITE, align=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.MIDDLE)
    body_y = y + HDR_H
    bh = body_height or Inches(estimate_body_inches(body_text, body_size, w))
    _add_text_box(slide, x, body_y, w, bh, body_text,
                  size=body_size, color=BLACK)
    return body_y + bh


def add_picture_block(slide, x, y, w, header, image_path, *,
                      header_size=32, max_h_in=None):
    HDR_H = Inches(1.10)
    _add_rect(slide, x, y, w, HDR_H, fill=MAROON)
    _add_text_box(slide, x, y, w, HDR_H, header,
                  font="Calibri", size=header_size, bold=True,
                  color=WHITE, align=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.MIDDLE)
    pic_y = y + HDR_H + Inches(0.05)
    if image_path and os.path.exists(image_path):
        from PIL import Image
        with Image.open(image_path) as im:
            iw, ih = im.size
        ar = ih / iw
        target_w_in = w / Inches(1)
        target_h_in = target_w_in * ar
        if max_h_in and target_h_in > max_h_in:
            target_h_in = max_h_in
            target_w_in = target_h_in / ar
        slide.shapes.add_picture(
            image_path,
            x + (w - Inches(target_w_in)) / 2,
            pic_y,
            width=Inches(target_w_in))
        return pic_y + Inches(target_h_in) + Inches(0.05)
    else:
        ph_h = Inches(max_h_in or 5)
        _add_rect(slide, x, pic_y, w, ph_h, fill=LIGHT_GRAY, line=GRAY)
        _add_text_box(slide, x, pic_y, w, ph_h,
                      f"[ Pending: {os.path.basename(image_path)} ]",
                      size=18, color=GRAY, align=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)
        return pic_y + ph_h + Inches(0.05)


# ---------------------------------------------------------------- main

def build():
    prs = Presentation()
    prs.slide_width = Inches(48)
    prs.slide_height = Inches(36)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ============================ TOP BANNER ============================
    BANNER_H = Inches(4.6)
    LOGO_W = Inches(4.5)        # reserved top-left for TAMU logo
    QR_W = Inches(4.0)          # reserved top-right for QR

    # Maroon fill across the full banner width (logo overlays it on user
    # side; we still draw to keep the title row visually unified)
    _add_rect(slide, 0, 0, prs.slide_width, BANNER_H, fill=MAROON)
    _add_rect(slide, 0, BANNER_H, prs.slide_width, Inches(0.30),
              fill=MAROON_DK)

    # Title block (between logo on left and QR on right)
    title_x = LOGO_W
    title_w = prs.slide_width - LOGO_W - QR_W

    _add_text_box(
        slide, title_x, Inches(0.40), title_w, Inches(2.20),
        "An Agentic AI Science Community for\n"
        "Automated Neural Operator Discovery",
        font="Calibri", size=66, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _add_text_box(
        slide, title_x, Inches(2.65), title_w, Inches(0.70),
        "Luis Loo    |    Advisor: Ulisses Braga-Neto    |    "
        "Department of Electrical & Computer Engineering, Texas A&M University",
        font="Calibri", size=24, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _add_text_box(
        slide, title_x, Inches(3.40), title_w, Inches(0.85),
        "Hot Topics Workshop on Agentic Scientific Computing and "
        "Scientific Machine Learning   |   ICERM, Brown University   |   "
        "May 9-10, 2026",
        font="Calibri", size=20, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Top-left logo placeholder (user keeps existing TAMU logo here)
    _add_text_box(
        slide, Inches(0.20), Inches(0.20), LOGO_W - Inches(0.20),
        BANNER_H - Inches(0.40),
        "[ TAMU LOGO ]",
        font="Calibri", size=14, color=WHITE, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)

    # Top-right QR code
    qr_path = os.path.join(FIG_DIR, "qr_repo.png")
    if os.path.exists(qr_path):
        qr_box_x = prs.slide_width - QR_W + Inches(0.30)
        qr_box_y = Inches(0.30)
        qr_box_h = BANNER_H - Inches(0.85)
        # white square behind the QR for readability
        _add_rect(slide, qr_box_x - Inches(0.10),
                  qr_box_y - Inches(0.10),
                  qr_box_h + Inches(0.20),
                  qr_box_h + Inches(0.20),
                  fill=WHITE)
        slide.shapes.add_picture(qr_path, qr_box_x, qr_box_y,
                                 height=qr_box_h)
        _add_text_box(
            slide, prs.slide_width - QR_W, BANNER_H - Inches(0.85),
            QR_W, Inches(0.70),
            "github.com/luislootx/AI-SC",
            font="Consolas", size=14, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ============================ BODY ============================
    body_top = Inches(5.20)
    margin = Inches(0.5)
    bottom_margin = Inches(0.5)
    body_h = prs.slide_height - body_top - bottom_margin

    # Left column = ~1/3 of width; right span = ~2/3 of width
    left_w = (prs.slide_width - 2 * margin) * 0.32
    gap = Inches(0.4)
    right_x = margin + left_w + gap
    right_w = prs.slide_width - margin - right_x

    # ============================ LEFT COLUMN ============================
    y = body_top
    y = add_section(slide, margin, y, left_w,
        "Key Takeaway",
        "A swarm of virtual labs, exchanging CITATIONS like a research "
        "community, discovers HYBRID neural operator architectures with "
        "no human prior - and converges on a DIFFERENT architecture for "
        "each PDE family.",
        body_size=26)

    y += Inches(0.30)
    y = add_section(slide, margin, y, left_w,
        "Motivation",
        "  - Neural operators (FNO, DeepONet, GNOT, MWT, ...) ship with "
        "fixed inductive biases. Picking the right family is still done "
        "by hand.\n\n"
        "  - QUESTION: can a swarm of virtual labs discover the right "
        "architecture per PDE, with no human prior?\n\n"
        "  - FRAMEWORK: Braga-Neto's AI Scientific Community "
        "(arXiv:2603.21344, 2026). 16 labs, citation-based influence, "
        "PSO coordinator over architecture genomes.",
        body_size=24)

    y += Inches(0.30)
    y = add_section(slide, margin, y, left_w,
        "Methods - Five Role Agents per Lab",
        "  - PLANNER:   proposes the architecture genome.\n"
        "  - WORKER:    inner GA over hyperparameters.\n"
        "  - EVALUATOR: multi-objective fitness\n"
        "      f = w_a acc + w_g gen + w_e eff + w_n nov + w_p peer.\n"
        "  - REVIEWER:  votes for top-2 peer labs.\n"
        "  - COORDINATOR (shared): PSO + culling.\n\n"
        "TRUST SCORE\n"
        "    tau_i = (1-a) tau_i + a c_i / max c.\n"
        "Drives GA budget, PSO weight, lab survival.\n\n"
        "BLOCK LIBRARY:  Fourier (FNO)  -  Branch-Trunk (DeepONet)\n"
        "  -  Wavelet (MWT)  -  Attention (GNOT)  -  Residual Conv.",
        body_size=22)

    y += Inches(0.30)
    y = add_section(slide, margin, y, left_w,
        "Discussion & Caveats",
        "  - Regime-aware hybrid: 4-family recombination, no human prior.\n"
        "  - For NS-2D the hybrid matches FNO h64 m12 with ~3x fewer params.\n"
        "  - Block usage shifts measurably iter 1 -> 20 in every benchmark - "
        "this is not random search.\n"
        "  - CAVEAT: absolute rel L2 is NOT comparable across PDEs (signal "
        "energies, noise floors, spectral structure differ).\n"
        "  - We do NOT claim to beat DeepONet. Darcy-2D test pending.",
        body_size=22)

    y += Inches(0.30)
    y = add_section(slide, margin, y, left_w,
        "Next Steps & References",
        "NEXT:  Darcy-2D (DeepONet regime); LLM-backed Planner / Reviewer "
        "(Athena, Toscano et al., 2025); graph operator + neural Galerkin "
        "blocks.\n\n"
        "  - Braga-Neto. AI Scientific Community. arXiv:2603.21344 (2026).\n"
        "  - Lu et al. The AI Scientist. Nature 651, 914-919 (2026).\n"
        "  - Toscano, Chen, Karniadakis. Athena. arXiv:2512.03476 (2025).\n"
        "  - Li et al. FNO. ICLR 2021.   Lu et al. DeepONet. NMI 3 (2021).\n\n"
        "ACKNOWLEDGMENTS  -  Prof. Ulisses Braga-Neto; ICERM organisers; "
        "TAMU ECEN travel grant; ICERM lodging support.",
        body_size=20)

    # ============================ RIGHT 2-COL SPAN ============================
    # Stack: Framework Overview, Master Results, Bottom row (architectures + table)
    y = body_top
    y = add_picture_block(slide, right_x, y, right_w,
        "Framework Overview",
        os.path.join(FIG_DIR, "architecture_overview.png"),
        max_h_in=11.0)

    y += Inches(0.30)
    y = add_picture_block(slide, right_x, y, right_w,
        "Results  -  All Four Experiments",
        os.path.join(FIG_DIR, "master_results.png"),
        max_h_in=12.0)

    # Bottom row: 2 side-by-side blocks (architectures + baselines table)
    y += Inches(0.30)
    bottom_w = (right_w - Inches(0.30)) / 2
    add_section(slide, right_x, y, bottom_w,
        "Discovered Architectures",
        "PIECEWISE REG. 1D       R + R + R\n"
        "LIN. ADVECTION 1D       A + F + F + T + A\n"
        "BURGERS 1D              F + R + F + T + A + F\n"
        "NS-2D (winner)          F + A + W + R + T + T + F\n"
        "NS-2D (best rel L2)     F + F + A + W + R + T + F\n"
        "                        rel L2 = 2.6e-4   (1.5M params)\n\n"
        "Different PDEs => different architectures.\n"
        "Legend: F=Fourier  A=Attention  W=Wavelet\n"
        "        R=ResConv  T=Branch-Trunk.",
        body_size=20)

    add_section(slide, right_x + bottom_w + Inches(0.30), y, bottom_w,
        "Honest Comparison  -  1D Baselines",
        "Same data, same training budget per architecture\n"
        "(40 epochs, 256 train, res 128).\n\n"
        "PDE          | Hybrid | FNO   | DeepO | Wave  | Trans\n"
        "Piecewise    | 0.043  | 0.062 | 0.224 | 0.043*| 0.078\n"
        "Lin. Advec.  | 0.020* | 0.033 | 0.866 | 0.546 | 0.862\n"
        "Burgers      | 0.074* | 0.117 | 0.791 | 0.430 | 0.784\n\n"
        "* = best per row.  Hybrid WINS Advection + Burgers,\n"
        "TIES the regime-optimal pure family on regression.",
        body_size=20)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
